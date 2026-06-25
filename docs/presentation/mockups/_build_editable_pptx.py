"""목업 추출 JSON(요소별 bbox/스타일)을 편집 가능한 PPTX 객체로 변환.
각 요소 = 도형(box)·텍스트박스(tx)·이미지(img). px→EMU(1px=9525), fs px→pt(*0.75).
폰트는 'Malgun Gothic'으로 통일(실제 Canva는 210 다락방/나눔스퀘어로 교체).
"""
import os, re, json
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

MK = os.path.dirname(os.path.abspath(__file__)).replace("\\", "/")
ALL = MK + "/_editable_data.json"  # 브라우저에서 추출한 요소별 bbox/스타일(이중 인코딩 가능)
OUT = MK + "/../no-home-mockups-editable.pptx"
EMU = 9525
FONT = "Malgun Gothic"

ORDER = [
    "11-problem", "12-required", "13-ai-intro", "21-milestone", "22-team",
    "31-competitor", "32-differentiation", "41-required-status", "42-search",
    "43-member", "44-ai-assistant", "45-verification", "51-architecture", "52-stack",
    "61-flow", "62-demo", "71-toolcalling", "72-capability", "73-memory",
    "74-performance", "75-security", "81-impact", "91-troubleshooting", "92-retro", "93-closing",
]
ALIGN = {"start": PP_ALIGN.LEFT, "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


def hexc(s):
    m = re.findall(r"\d+", s or "")
    return RGBColor(int(m[0]), int(m[1]), int(m[2])) if len(m) >= 3 else RGBColor(0, 0, 0)


def E(px):
    return Emu(int(round(px * EMU)))


def build(prs, elems):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    for el in elems:
        t = el["t"]
        if t == "box":
            w, h, radius = el["w"], el["h"], el.get("radius", 0)
            if radius and radius >= min(w, h) / 2:
                shp = MSO_SHAPE.OVAL
            elif radius:
                shp = MSO_SHAPE.ROUNDED_RECTANGLE
            else:
                shp = MSO_SHAPE.RECTANGLE
            sp = s.shapes.add_shape(shp, E(el["x"]), E(el["y"]), E(w), E(h))
            sp.line.fill.background()
            sp.shadow.inherit = False
            if el.get("bg"):
                sp.fill.solid(); sp.fill.fore_color.rgb = hexc(el["bg"])
            else:
                sp.fill.background()
            if el.get("blw", 0) >= 3 and el.get("blc"):
                bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(el["x"]), E(el["y"]), E(el["blw"]), E(h))
                bar.line.fill.background(); bar.shadow.inherit = False
                bar.fill.solid(); bar.fill.fore_color.rgb = hexc(el["blc"])
        elif t == "img":
            path = os.path.normpath(os.path.join(MK, el["src"]))
            if os.path.exists(path):
                s.shapes.add_picture(path, E(el["x"]), E(el["y"]), E(el["w"]), E(el["h"]))
        elif t == "tx":
            tb = s.shapes.add_textbox(E(el["x"]), E(el["y"]), E(el["w"]), E(el["h"]))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.alignment = ALIGN.get(el.get("al"), PP_ALIGN.LEFT)
            r = p.add_run(); r.text = el["s"]
            f = r.font
            f.size = Pt(round(el["fs"] * 0.75))
            f.bold = el.get("fw", 400) >= 700
            f.color.rgb = hexc(el["col"])
            f.name = FONT
    return s


prs = Presentation()
prs.slide_width = Emu(int(1280 * EMU))
prs.slide_height = Emu(int(720 * EMU))

_raw = open(ALL, encoding="utf-8").read()
data = json.loads(_raw)
if isinstance(data, str):
    data = json.loads(data)

n = 0
for f in ORDER:
    elems = data.get(f)
    if not elems:
        print("MISSING/EMPTY", f); continue
    build(prs, elems)
    n += 1

prs.save(OUT)
print(f"SAVED {OUT} ({n} slides)")

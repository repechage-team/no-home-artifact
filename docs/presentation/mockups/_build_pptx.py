"""목업 HTML 캡처(raw PNG)를 crop 후 16:9 PPTX로 full-bleed 조립.
캡처: Chrome headless --window-size=1328,768 --force-device-scale-factor=2 (slide 1280x720, body padding 24)
crop : (48,48)~(2608,1488) = 슬라이드 1280x720 @2x = 2560x1440
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

MK = os.path.dirname(os.path.abspath(__file__)).replace("\\", "/")
RAW = MK + "/_pptx_build/raw"
CROP = MK + "/_pptx_build/crop"
OUT = MK + "/../no-home-mockups.pptx"  # docs/presentation/no-home-mockups.pptx
os.makedirs(CROP, exist_ok=True)

ORDER = [
    "11-problem", "12-required", "13-ai-intro",
    "21-milestone", "22-team",
    "31-competitor", "32-differentiation",
    "41-required-status", "42-search", "43-member", "44-ai-assistant", "45-verification",
    "51-architecture", "52-stack",
    "61-flow", "62-demo",
    "71-toolcalling", "72-capability", "73-memory", "74-performance", "75-security",
    "81-impact",
    "91-troubleshooting", "92-retro", "93-closing",
]
BOX = (48, 48, 48 + 2560, 48 + 1440)  # slide area @2x

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

n = 0
for f in ORDER:
    raw = f"{RAW}/{f}.png"
    if not os.path.exists(raw):
        print("MISSING", f); continue
    im = Image.open(raw).crop(BOX)
    cp = f"{CROP}/{f}.png"
    im.save(cp)
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(cp, 0, 0, width=prs.slide_width, height=prs.slide_height)
    n += 1

prs.save(OUT)
print(f"SAVED {OUT} ({n} slides)")
